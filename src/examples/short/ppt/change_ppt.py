"""
Use the python-ppt to change a ppt

References:
- https://python-pptx.readthedocs.io/en/latest/#user-guide
"""
from pptx import Presentation

input_filename = "data_samples/sample.pptx"
output_filename = "/tmp/out.pptx"

presentation = Presentation(input_filename)
slides = presentation.slides
for slide in slides:
    print(dir(slide))
    print(slide.follow_master_background)
    # slide.follow_master_background = False
    # if slide.background is None:
    #     print("background is none")
    #     continue
    print(len(slide.shapes))
    # for shape in slide.shapes:
    # print(type(shape.shape_type))
    # print(dir(shape.shape_type))
    #     if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
    #        print("deleting")
    # shape.remove(shape)
    #         slide.shapes.remove(shape)
    # slide.remove(shape)
    element = slide.background.element
    element.delete()
    # slide.background.element = None
    # print(dir(element))
    # slide.background = None
# print(dir(presentation))
presentation.save(output_filename)
