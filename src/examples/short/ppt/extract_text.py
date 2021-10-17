"""
Extract text from a ppt file

References:
- https://www.syncfusion.com/kb/8672/how-to-extract-text-from-a-powerpoint-presentation
- https://python-pptx.readthedocs.io/en/latest/#user-guide
"""
import sys
from itertools import islice

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

input_filename = sys.argv[1]

presentation = Presentation(input_filename)
slides = presentation.slides
for slide_number, slide in enumerate(islice(slides, None)):
    print(f"slide number {slide_number}")
    if slide.name != "":
        print(f"slide.name {slide.name}")
    # This block extracts hyperlinks
    for v in slide.part.rels.values():
        target: str = v.target_ref
        if target.startswith(".."):
            continue
        print(v.target_ref)
    for shape in slide.shapes:
        # pylint: disable=no-member
        if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
            print(shape.shape_type)
            print(shape.text)
        if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
            print(shape.shape_type)
            print(shape.text)
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            print(shape.shape_type)
            for row in shape.table.rows:
                for cell in row.cells:
                    print(cell.text)
